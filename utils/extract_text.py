import PyPDF2
from docx import Document
from pptx import Presentation


def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF file."""
    text = ""
    with open(pdf_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() or ""
    return text


def extract_text_from_word(word_path):
    """Extract text from a Word (.docx) file."""
    doc = Document(word_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text


def extract_text_from_powerpoint(ppt_path):
    """Extract text from a PowerPoint (.pptx) file."""
    prs = Presentation(ppt_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

if __name__ == "__main__":
    pdf_path = "data/Chain_of_thought.pdf"
    word_path = "data/LỊCH TRỰC PHÒNG.docx"
    ppt_path = "data/Final Report TMA.pptx"

    text = extract_text_from_word(word_path)
    print(text)

