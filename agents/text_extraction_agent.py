import os
import sys
import PyPDF2
from docx import Document
from pptx import Presentation

# Thêm đường dẫn cha để import module cấu hình model
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from langchain_core.tools import tool
from langchain.prompts import PromptTemplate
from langgraph.prebuilt import create_react_agent
from config.llm import gemini

# ---------------------- TOOLS ----------------------

@tool("extract_text_from_pdf")
def extract_text_from_pdf(pdf_path: str) -> str:
    """Sử dụng công cụ này để trích xuất văn bản từ tệp PDF. Đầu vào là một chuỗi đường dẫn đến tệp .pdf."""
    text = ""
    try:
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() or ""
        return text if text else "Không tìm thấy văn bản có thể trích xuất trong PDF."
    except Exception as e:
        return f"Lỗi khi trích xuất PDF: {e}"

@tool("extract_text_from_word")
def extract_text_from_word(word_path: str) -> str:
    """Sử dụng công cụ này để trích xuất văn bản từ tệp Word (.docx). Đầu vào là một chuỗi đường dẫn đến tệp .docx."""
    try:
        doc = Document(word_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text if text else "Không tìm thấy văn bản trong tài liệu Word."
    except Exception as e:
        return f"Lỗi khi trích xuất Word: {e}"

@tool("extract_text_from_powerpoint")
def extract_text_from_powerpoint(ppt_path: str) -> str:
    """Sử dụng công cụ này để trích xuất văn bản từ tệp PowerPoint (.pptx). Đầu vào là một chuỗi đường dẫn đến tệp .pptx."""
    try:
        prs = Presentation(ppt_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text if text else "Không tìm thấy văn bản trong các slide PowerPoint."
    except Exception as e:
        return f"Lỗi khi trích xuất PowerPoint: {e}"

# ---------------------- AGENT SETUP ----------------------

def create_text_extraction_agent():
    """Tạo một ReAct Agent để trích xuất văn bản."""
    tools = [
        extract_text_from_pdf,
        extract_text_from_word,
        extract_text_from_powerpoint
    ]

    text_extraction_prompt = """
### Vai trò:
Bạn là **một trợ lý tài liệu thông minh**, chuyên xử lý và **trích xuất nội dung văn bản từ các loại tài liệu** khác nhau.

---

### Mục tiêu:
Trích xuất **nội dung văn bản thuần túy** từ các tệp tài liệu sau, dựa trên phần mở rộng của tệp:

- **PDF** (`.pdf`) → sử dụng công cụ `[extract_text_from_pdf]`  
- **Microsoft Word** (`.docx`) → sử dụng công cụ `[extract_text_from_word]`  
- **Microsoft PowerPoint** (`.pptx`) → sử dụng công cụ `[extract_text_from_powerpoint]`  

---

### Hành vi mong muốn:
- Tự động chọn đúng công cụ dựa trên phần mở rộng của đường dẫn `file_path`.
- Mỗi công cụ chỉ nhận một đối số: đường dẫn tuyệt đối tới tệp tài liệu (`file_path`).
- **Không xử lý định dạng** (không cần bảng, hình ảnh, biểu đồ...).
- **Chỉ trả về phần văn bản đã trích xuất**, không thêm mô tả, tiêu đề hay nhận xét.

---

### Đầu ra:
> Một chuỗi chứa **văn bản thuần túy (`plain text`)** được trích xuất từ tài liệu đầu vào.

---

### Lưu ý:
- Nếu tài liệu không thể trích xuất được văn bản, trả về đúng dòng:  
  `"Không thể trích xuất nội dung từ tệp này."`
"""

    agent = create_react_agent(
        tools=tools,
        model=gemini,
        prompt=text_extraction_prompt,
        name="Text Extraction Agent"
    )
    return agent

# ---------------------- RUN TEST ----------------------

if __name__ == "__main__":
    agent = create_text_extraction_agent()
    result = agent.invoke(
        {
            "messages": "Trích xuất văn bản từ tệp tại đường dẫn: D:/Project/Chatbot_CNM/data/Chain_of_thought.pdf"
        },
        config={"recursion_limit": 50}
    )

    # In toàn bộ message để debug nếu cần
    for msg in result["messages"]:
        print(msg.content)